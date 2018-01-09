import zipfile
import os
import json
import sys
import win32com
import pptx
import shutil
from pptx import Presentation
import ImageSerialization
import DatabaseConnection

def PPT_extract_caption(filepath):

    presentation = Presentation(filepath)
    titles = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if 'Picture' in str(shape.name):
                if slide.shapes.title:
                    titles.append(slide.shapes.title.text_frame.text)
                else:
                    titles.append(os.path.splitext(os.path.basename(filepath))[0])

    return titles

def saveImage(path, picName):

    pic = os.path.basename(picName)

    if not os.path.exists("result"):
        os.mkdir("result")

    newpath = os.path.join("result", pic)

    fout = open(newpath, 'w')
    fout.close()

    with open(path, 'rb') as fin:
        with open(newpath, 'ab') as fout:
            buff = fin.read(128)
            while buff:
                fout.write(buff)
                buff = fin.read(128)

def PPT_convertor(filepath):
    try:
        Application = win32com.client.Dispatch("PowerPoint.Application")
        Application.Visible = True
        if filepath.endswith(".ppt"):
            new_filepath = filepath.replace(".ppt", r".pptx")
        elif filepath.endswith(".pps"):
            new_filepath = filepath.replace(".pps", r".pptx")
        Presentation = Application.Presentations.Open(filepath)
        Presentation.Saveas(new_filepath)
        Presentation.Close()
        return new_filepath
    except Exception as e:
        print(e)
    finally:
        Application.Quit()

def PPTX_extractImages(filepath):

    if zipfile.is_zipfile(filepath):

        z = zipfile.ZipFile(filepath)
        images = list()
        for file in z.namelist():
            path, name = os.path.split(file)
            if path == "ppt/media":
                images.append(file)

        caption = PPT_extract_caption(filepath)

        finalResult = []
        contor = 0
        for img in images:
            z.extract(img, r'.')
            contor += 1
            result_json = ImageSerialization.get_info_from_image(img)

            newpath = os.path.join('.', img)
            saveImage(newpath,img)

            if contor < len(caption):
                result_json['caption'] = caption[contor]
            else:
                result_json['caption'] = os.path.splitext(os.path.basename(filepath))[0]

            del result_json['pixels']
            result_json['path'] = os.path.join('.','result', os.path.basename(img))
            finalResult.append(result_json)
            contor += 1

        file = open("%sresult.json" % argv[0], 'w')
        json.dump(finalResult, file, separators=(',', ':'))

        db = DatabaseConnection.DatabaseConnection()
        db.insert_entry(finalResult)

        file.close()
        shutil.rmtree('./ppt')


def ppt_document_parser(argv):

    if len(argv) != 2:
        print("Wrong number of parameters")
        exit(0)

    filepath = argv[1]

    if filepath.endswith(".pptx"):
        PPTX_extractImages(filepath)

    elif filepath.endswith(".ppt") or filepath.endswith(".pps"):

        new_filepath = PPT_convertor(filepath)
        PPTX_extractImages(new_filepath)
        os.remove(new_filepath)

if __name__ == '__main__':
    argv = sys.argv
    ppt_document_parser(argv)